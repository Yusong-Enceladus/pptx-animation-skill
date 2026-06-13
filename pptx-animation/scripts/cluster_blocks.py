#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Auto-cluster every slide's shapes into SEMANTIC BLOCKS and emit an animation spec for
inject_animations.py. A "block" is a whole visual unit — a card, a column, a row, a table, a
panel-with-its-rows — that should fade in together on ONE click (NOT element-by-element).

This encodes the hard-won lesson that good click animation is BLOCK-level and ordered by reading
logic. The algorithm:
  1. drop chrome (header band/title/page number) by position; separate arrow autoshapes.
  2. pull thin full-width bars (lead text / sub-headers / banners) so they don't block body cuts.
  3. GUILLOTINE the body: recursively split along the largest clean gap (>GAP) — background panels
     keep a card/column/row intact; the gaps BETWEEN units become the cuts. Adapts per slide:
     column layouts cut into columns, row layouts (each row has a full-width bg) cut into rows,
     dense tables don't cut (one block).
  4. merge passes that fix the residue:
       - list-collapse: a run of >=2 wide-thin rows -> one list block (a citation list).
       - containment: a block that mostly overlaps another (panel behind a list, lead over a flow
         chain) merges in -> no empty background-panel block, no orphaned pills.
       - vertical: stack column-internal pieces (header+body+image of one narrow column).
       - short-row / narrow-row: chip rows, bottom flow steps, cert/photo/tile rows.
     A `separated()` guard stops a vertical merge across a full-width separator bar.
  5. attach arrows to the nearest block; order blocks by READING ORDER — group into rows by real
     y-overlap (side-by-side only; stacked blocks go to separate rows), rows top->bottom, within a
     row left->right; a top lead bar is forced first, a full-width bottom banner forced last.

Render a numbered block-map (--map dir/) to verify the decomposition WITHOUT opening PowerPoint.

Usage:
  python cluster_blocks.py deck.pptx --out spec.json [--skip 1,2,3] [--map mapdir] [--duration 0.45]
"""
import sys, json, argparse, re
from pptx import Presentation
EMU=914400.0
GAP=0.16; SPAN=8.5

def is_arrow(sh):
    try: n=str(sh.auto_shape_type); return 'ARROW' in n or 'CHEVRON' in n
    except Exception: return False
def xov(a,b): return max(0,min(a['x2'],b['x2'])-max(a['x1'],b['x1']))
def bbox(blk): return dict(x1=min(s['x1'] for s in blk),y1=min(s['y1'] for s in blk),
                           x2=max(s['x2'] for s in blk),y2=max(s['y2'] for s in blk))

def guillotine(shapes):
    if len(shapes)<=1: return [shapes]
    best=None
    for ax in ('x','y'):
        iv=sorted((s[ax+'1'],s[ax+'2']) for s in shapes); maxend=iv[0][1]
        for a,b in iv[1:]:
            g=a-maxend
            if g>GAP and (best is None or g>best[0]): best=(g,ax,maxend+g/2)
            if b>maxend: maxend=b
    if best is None: return [shapes]
    _,ax,pos=best
    lo=[s for s in shapes if (s[ax+'1']+s[ax+'2'])/2<pos]; hi=[s for s in shapes if (s[ax+'1']+s[ax+'2'])/2>=pos]
    return guillotine(lo)+guillotine(hi)

def cluster(slide, header_h=1.05):
    shapes=[]; arrows=[]
    for sh in slide.shapes:
        x,y,w,h=sh.left,sh.top,sh.width,sh.height
        if None in (x,y,w,h): continue
        if w>=13.0*EMU and h>=7.0*EMU: continue                 # full-slide background
        if y<header_h*EMU: continue                             # header band / title / tag
        if y>6.3*EMU and x>11.0*EMU and w<1.2*EMU: continue     # page number (bottom-right)
        if y>6.85*EMU and h<0.45*EMU: continue                  # bottom source/footnote line
        r=dict(id=sh.shape_id,x1=x/EMU,y1=y/EMU,x2=(x+w)/EMU,y2=(y+h)/EMU)
        (arrows if is_arrow(sh) else shapes).append(r)
    if len(shapes)<2: return []
    spanners=[s for s in shapes if (s['x2']-s['x1'])>SPAN and (s['y2']-s['y1'])<0.85]
    body=[s for s in shapes if not ((s['x2']-s['x1'])>SPAN and (s['y2']-s['y1'])<0.85)]
    blocks=[g for g in guillotine(body) if g] if body else []
    blocks+=[[s] for s in spanners]
    bb=bbox
    def separated(a,b):
        ylo=min(a['y2'],b['y2']); yhi=max(a['y1'],b['y1'])
        for s in spanners:
            scy=(s['y1']+s['y2'])/2
            if ylo-0.06<=scy<=yhi+0.06 and s['x1']<=max(a['x1'],b['x1'])+0.1 and s['x2']>=min(a['x2'],b['x2'])-0.1: return True
        return False
    def domerge(test):
        m=True
        while m:
            m=False
            for i in range(len(blocks)):
                for j in range(i+1,len(blocks)):
                    if test(bb(blocks[i]),bb(blocks[j])): blocks[i]=blocks[i]+blocks[j]; blocks.pop(j); m=True; break
                if m: break
    # list-collapse (runs of wide-thin rows -> one list)
    is_row=lambda blk:(bb(blk)['x2']-bb(blk)['x1'])>6 and (bb(blk)['y2']-bb(blk)['y1'])<0.75
    thin=sorted([i for i in range(len(blocks)) if is_row(blocks[i])], key=lambda i:bb(blocks[i])['y1'])
    runs=[]; cur=[]
    for i in thin:
        if cur:
            p=bb(blocks[cur[-1]]); c=bb(blocks[i])
            if c['y1']-p['y2']<0.45 and xov(p,c)>0.5*(c['x2']-c['x1']) and not separated(p,c): cur.append(i); continue
            if len(cur)>=2: runs.append(cur)
            cur=[i]
        else: cur=[i]
    if len(cur)>=2: runs.append(cur)
    rm=set()
    for run in runs:
        for j in run[1:]: blocks[run[0]]=blocks[run[0]]+blocks[j]; rm.add(j)
    blocks=[b for k,b in enumerate(blocks) if k not in rm]
    # containment (mostly-overlap)
    def ctest(a,b):
        ox=max(0,min(a['x2'],b['x2'])-max(a['x1'],b['x1'])); oy=max(0,min(a['y2'],b['y2'])-max(a['y1'],b['y1']))
        aA=(a['x2']-a['x1'])*(a['y2']-a['y1']); bA=(b['x2']-b['x1'])*(b['y2']-b['y1'])
        return ox*oy>0.6*min(aA,bA)
    domerge(ctest)
    # vertical column-internal
    def vtest(a,b):
        wa=a['x2']-a['x1']; wb=b['x2']-b['x1']; ov=xov(a,b); ygap=max(b['y1']-a['y2'],a['y1']-b['y2'])
        return wa<8 and wb<8 and ov>0.5*min(wa,wb) and ygap<0.8 and not separated(a,b)
    domerge(vtest)
    # short same-band (chip rows / flow steps)
    def stest(a,b):
        ha=a['y2']-a['y1']; hb=b['y2']-b['y1']; band=abs((a['y1']+a['y2'])/2-(b['y1']+b['y2'])/2)<0.5
        xgap=max(b['x1']-a['x2'],a['x1']-b['x2']); return ha<1.0 and hb<1.0 and band and xgap<1.3
    domerge(stest)
    # narrow same-band (cert/photo/tile rows)
    def ntest(a,b):
        wa=a['x2']-a['x1']; wb=b['x2']-b['x1']; band=abs((a['y1']+a['y2'])/2-(b['y1']+b['y2'])/2)<0.6
        xgap=max(b['x1']-a['x2'],a['x1']-b['x2']); return wa<2.5 and wb<2.5 and band and xgap<0.6
    domerge(ntest)
    # attach arrows to nearest block
    for ar in arrows:
        best=None; bd=9e9
        for k,blk in enumerate(blocks):
            c=bb(blk); d=(((c['x1']+c['x2'])/2-(ar['x1']+ar['x2'])/2)**2+((c['y1']+c['y2'])/2-(ar['y1']+ar['y2'])/2)**2)
            if d<bd: bd=d; best=k
        if best is not None: blocks[best].append(ar)
    # reading order
    bbs=[bb(b) for b in blocks]
    is_lead=lambda b:(b['x2']-b['x1'])>7.5 and (b['y2']-b['y1'])<0.95 and (b['y1']+b['y2'])/2<1.85
    is_ban =lambda b:(b['x2']-b['x1'])>7.5 and (b['y2']-b['y1'])<0.95 and (b['y1']+b['y2'])/2>5.5
    leads=[i for i in range(len(blocks)) if is_lead(bbs[i])]; bans=[i for i in range(len(blocks)) if is_ban(bbs[i])]
    mid=[i for i in range(len(blocks)) if i not in leads and i not in bans]; mid.sort(key=lambda i:bbs[i]['y1'])
    rows=[]
    for i in mid:
        placed=False
        for r in rows:
            yov=min(bbs[i]['y2'],r['y2'])-max(bbs[i]['y1'],r['y1'])
            if yov<=0.4*min(bbs[i]['y2']-bbs[i]['y1'],r['y2']-r['y1']): continue
            if any(xov(bbs[i],bbs[j])>0.5*min(bbs[i]['x2']-bbs[i]['x1'],bbs[j]['x2']-bbs[j]['x1']) for j in r['ix']): continue
            r['ix'].append(i); r['y1']=min(r['y1'],bbs[i]['y1']); r['y2']=max(r['y2'],bbs[i]['y2']); placed=True; break
        if not placed: rows.append({'y1':bbs[i]['y1'],'y2':bbs[i]['y2'],'ix':[i]})
    rows.sort(key=lambda r:r['y1'])
    order=sorted(leads,key=lambda i:bbs[i]['y1'])
    for r in rows: order+=sorted(r['ix'],key=lambda i:((bbs[i]['x1']+bbs[i]['x2'])/2,bbs[i]['y1']))
    order+=sorted(bans,key=lambda i:bbs[i]['y1'])
    return [[s['id'] for s in blocks[i]] for i in order], [ (bbs[i]['x1'],bbs[i]['y1'],bbs[i]['x2'],bbs[i]['y2']) for i in order ]

def main():
    ap=argparse.ArgumentParser()
    ap.add_argument('pptx'); ap.add_argument('--out',required=True); ap.add_argument('--duration',type=float,default=0.45)
    ap.add_argument('--skip',default='',help='comma-separated 1-based presentation indices to skip (cover/dividers/spliced/etc.)')
    ap.add_argument('--map',default=None,help='dir to write numbered block-map JPGs (needs the deck rendered as <dir>/p-NN.jpg at 96dpi)')
    a=ap.parse_args()
    skip=set(int(x) for x in a.skip.split(',') if x.strip())
    prs=Presentation(a.pptx); spec={'duration':a.duration,'slides':[]}; maps={}
    for idx,slide in enumerate(prs.slides,start=1):
        if idx in skip: continue
        res=cluster(slide)
        if not res: continue
        groups,boxes=res
        if not groups: continue
        spec['slides'].append({'slide':idx,'groups':groups}); maps[idx]=boxes
    json.dump(spec,open(a.out,'w'),ensure_ascii=False,indent=1)
    clicks=sum(len(s['groups']) for s in spec['slides'])
    print(f"{len(spec['slides'])} slides, {clicks} blocks (avg {clicks/max(1,len(spec['slides'])):.1f}/slide) -> {a.out}")
    if a.map:
        json.dump(maps,open(a.map.rstrip('/')+'/blockmap.json','w'))
        print(f"block boxes -> {a.map}/blockmap.json (overlay on 96dpi renders to review)")

if __name__=='__main__': main()
